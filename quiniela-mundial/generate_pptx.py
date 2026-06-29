from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    two_content_layout = prs.slide_layouts[3]
    
    # --- SLIDE 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Quiniela Mundial 2026"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Una experiencia interactiva y autónoma para vivir la Copa del Mundo\n(Sin dolores de cabeza operativos)"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # --- SLIDE 2: Resumen Ejecutivo ---
    slide = prs.slides.add_slide(title_and_content_layout)
    title = slide.shapes.title
    title.text = "Resumen Ejecutivo"
    
    content = slide.placeholders[1].text_frame
    content.text = "La Quiniela Mundial 2026 es una PWA (aplicación web progresiva) de alto rendimiento diseñada para la competencia en tiempo real entre amigos y compañeros."
    p = content.add_paragraph()
    p.text = "Destaca por:"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Sincronización autónoma de marcadores (vía ESPN)"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Motor inteligente de puntuación (+3 pts exacto, +1 pt tendencia)"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Gestión 100% en la Nube (Firebase) sin necesidad de tocar código fuente"
    p.level = 1
    
    # --- SLIDE 3: Arquitectura y Cero Mantenimiento ---
    slide = prs.slides.add_slide(two_content_layout)
    title = slide.shapes.title
    title.text = "Arquitectura 100% Cloud: Cero Mantenimiento"
    
    left_body = slide.placeholders[1].text_frame
    left_body.text = "Tecnología:"
    p = left_body.add_paragraph()
    p.text = "Frontend Moderno: HTML/CSS/JS (PWA)"
    p.level = 1
    p = left_body.add_paragraph()
    p.text = "Backend DB: Firebase Firestore"
    p.level = 1
    p = left_body.add_paragraph()
    p.text = "Scripts Python: Control total de base de datos"
    p.level = 1
    p = left_body.add_paragraph()
    p.text = "BigQuery: Analítica avanzada sincronizada"
    p.level = 1
    
    right_body = slide.placeholders[2].text_frame
    right_body.text = "El Valor Diferencial:"
    p = right_body.add_paragraph()
    p.text = "La jerarquía de datos ha sido invertida. Firebase es el Jefe Absoluto."
    p.level = 1
    p = right_body.add_paragraph()
    p.text = "Si cambia un horario, estadio o probabilidad, el administrador simplemente corre un script en Python en su máquina local."
    p.level = 1
    p = right_body.add_paragraph()
    p.text = "Todos los usuarios de la web ven la actualización de inmediato. ¡Adiós deploys engorrosos en Netlify!"
    p.level = 1
    
    # --- SLIDE 4: Espacio para Screenshots ---
    slide = prs.slides.add_slide(title_and_content_layout)
    title = slide.shapes.title
    title.text = "Experiencia de Usuario (UI/UX)"
    content = slide.placeholders[1].text_frame
    content.text = "[ Por favor, inserta aquí las capturas de pantalla reales de tu aplicación ]\n\nSugerencias de capturas:\n1. Dashboard de Predicciones con el marcador de puntos sumado correctamente.\n2. La tarjeta de partido en vivo con la etiqueta de ESPN Live sincronizada."
    content.paragraphs[0].font.italic = True
    content.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    # Save the presentation to the Desktop
    desktop_path = os.path.join(os.path.expanduser("~"), "Desktop", "Quiniela_Mundial_OnePager.pptx")
    prs.save(desktop_path)
    print(f"Presentation saved to {desktop_path}")

if __name__ == "__main__":
    create_presentation()
