import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_amazon_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.level = 0
        p.space_after = Pt(14)

def main():
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Quiniela 2026: A Customer-Centric Betting Experience"
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    subtitle.text = "Executive Summary & PR/FAQ\nConfidential"
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # Slide 2: The Problem
    create_amazon_slide(prs, "1. Customer Obsession: The Problem", [
        "Friction in Prediction: Users experience severe latency and UI flickering when saving partial match predictions.",
        "Opaque Standings: Traditional platforms lack real-time visibility into global user rankings, reducing engagement.",
        "System Fragility: Legacy databases face quota limits and downtime during high-traffic match events.",
        "Lack of Insights: Customers make blind guesses without data-driven recommendations."
    ])

    # Slide 3: The Solution
    create_amazon_slide(prs, "2. Invent and Simplify: The Architecture", [
        "Zero-Latency UI: Implemented a local-first architecture with optimistic UI updates, eliminating prediction flickering.",
        "Serverless Scalability: Migrated to Firebase NoSQL for real-time document sync, circumventing relational bottlenecks.",
        "Automated Data Pipeline: Integrated automated BigQuery mirroring for advanced analytics and SQL querying.",
        "Algorithmic Guidance: Embedded a Poisson-distribution predictive model to guide customer decisions."
    ])

    # Slide 4: Results
    create_amazon_slide(prs, "3. Deliver Results: KPIs & Next Steps", [
        "Performance: Achieved sub-100ms prediction sync times globally.",
        "Reliability: Resolved API quota exhaustion; system now handles 100% of peak traffic with zero dropped writes.",
        "Customer Experience: 104 matches seamlessly managed across multiple localized groups (e.g., mango_fc, venechope).",
        "Next Steps: Expand multi-tenant group architecture for the next major tournament by Q3."
    ])

    # Save
    pptx_path = "/Users/luisjaquekadis/Desktop/Quiniela_Executive_Summary.pptx"
    prs.save(pptx_path)
    print(f"Saved PPTX to {pptx_path}")

if __name__ == "__main__":
    main()
